"""
Document Processor
Extracts text from various document formats
"""

import os
import time
from typing import List, Dict
from pathlib import Path

# PDF
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# DOCX
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Markdown
try:
    import markdown
    from bs4 import BeautifulSoup
    MD_AVAILABLE = True
except ImportError:
    MD_AVAILABLE = False

class DocumentProcessor:
    def __init__(self, chunk_size: int = 512):
        self.chunk_size = chunk_size  # tokens (~2000 characters)
        self.overlap = 50  # token overlap

    async def process_document(self, filepath: str) -> List[Dict]:
        """Process document and return text chunks"""
        file_ext = Path(filepath).suffix.lower()

        # Extract text based on file type
        if file_ext == '.pdf' and PDF_AVAILABLE:
            text = self._extract_pdf(filepath)
        elif file_ext == '.docx' and DOCX_AVAILABLE:
            text = self._extract_docx(filepath)
        elif file_ext == '.pptx' and PPTX_AVAILABLE:
            text = self._extract_pptx(filepath)
        elif file_ext in ['.md', '.markdown'] and MD_AVAILABLE:
            text = self._extract_markdown(filepath)
        elif file_ext == '.txt':
            text = self._extract_txt(filepath)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

        # Chunk text
        chunks = self._chunk_text(text)

        # Create chunk objects
        filename = Path(filepath).name
        result = []
        for i, chunk in enumerate(chunks):
            result.append({
                "text": chunk,
                "filename": filename,
                "filepath": filepath,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "file_type": file_ext[1:],
                "indexed_at": int(time.time())
            })

        return result

    def _extract_pdf(self, filepath: str) -> str:
        """Extract text from PDF"""
        text = []
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)

    def _extract_docx(self, filepath: str) -> str:
        """Extract text from DOCX"""
        doc = Document(filepath)
        return '\n'.join([para.text for para in doc.paragraphs])

    def _extract_pptx(self, filepath: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    def _extract_markdown(self, filepath: str) -> str:
        """Extract text from Markdown"""
        with open(filepath, 'r', encoding='utf-8') as f:
            md_text = f.read()
        # Convert MD to HTML, then extract text
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')
        return soup.get_text()

    def _extract_txt(self, filepath: str) -> str:
        """Extract text from TXT"""
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

    def _chunk_text(self, text: str) -> List[str]:
        """Chunk text with overlap"""
        # Simple chunking by characters (approx 4 chars = 1 token)
        chunk_chars = self.chunk_size * 4
        overlap_chars = self.overlap * 4

        chunks = []
        start = 0
        while start < len(text):
            end = start + chunk_chars
            chunk = text[start:end]

            # Try to break at sentence boundary
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)
                if break_point > chunk_chars * 0.5:  # At least 50% of chunk
                    chunk = chunk[:break_point + 1]
                    end = start + break_point + 1

            chunks.append(chunk.strip())
            start = end - overlap_chars if end < len(text) else end

        return [c for c in chunks if c]  # Remove empty chunks
